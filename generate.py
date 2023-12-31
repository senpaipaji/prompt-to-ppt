from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from langchain.llms import HuggingFaceHub
from langchain.chains import LLMChain
from langchain.prompts import PromptTemplate

def CreateChain(prompt,model='tiiuae/falcon-7b-instruct',length=512):
    print("Creating Chain!!")
    #model
    hub_llm = HuggingFaceHub(repo_id=model)
    #chain
    print("Finished Creating Chain!!")
    return LLMChain(prompt=prompt,llm=hub_llm,verbose=True)

def CreateTopics(topic):
    errors = []
    print("Creating Topics!!")
    try:
        prompt = PromptTemplate(
            input_variables = ["question"],
            template ="Without indices give only title of important subtopics of the topic {question}\n"
        )
        Extract = CreateChain(prompt).run(topic)
        Extract = Extract.split('\n')
        return Extract
    
    except Exception as e:
        errors.append(e)
        print("Finished Creating Topics!!")
        print(errors)

def CreateData(sections):
    errors = []
    print("Creating Data!!")
    presentation_data = {}
    prompt = PromptTemplate(
                input_variables = ["question"],
                template ="Explain in a single paragraph {question} \n"
            )
    Chain = CreateChain(prompt,length=100)
    for section in sections:
        try:
            title = section
            content = Chain.run(title)
            presentation_data[title] = content #saving
        except Exception as e:
            errors.append(e)
            continue
    if errors:
        print(errors)
    print("Finished Creating Data!!")
    return presentation_data

def TransformData(presentation_data):
    print('Performing Transformation!!')
    for title,content in presentation_data.items():
        temp_arr = []
        for item in presentation_data[title].split("\n\n"):
            temp_arr.append(item)
        presentation_data[title] = temp_arr
    print('Transformation Complete!!')
    return presentation_data

class PPT:
    def __init__(self, title, subtitle="  "):
        self.presentation = Presentation()
        self.title_slide_layout = self.presentation.slide_layouts[0]  
        self.content_slide_layout = self.presentation.slide_layouts[1]  
        
        # Title slide creation
        title_slide = self.presentation.slides.add_slide(self.title_slide_layout) 
        title_slide_title = title_slide.shapes.title
        title_slide_subtitle = title_slide.placeholders[1]
        title_slide_title.text = str(title)
        title_slide_subtitle.text = str(subtitle)
        print('PPT object created!!')
        
        
    def setTheme(self, theme_number = 0):
        match theme_number:
            case 0:
                dark1 = RGBColor(0, 0, 0)
                light1 = RGBColor(255, 255, 255)
                dark2 = RGBColor(30, 81, 85)
                light2 = RGBColor(235, 235, 235)

                accent1 = RGBColor(176, 21, 19)
                accent2 = RGBColor(234, 99, 18)
                accent3 = RGBColor(230, 183, 41)
                accent4 = RGBColor(106, 172, 144)
                accent5 = RGBColor(95, 156, 157)
                accent6 = RGBColor(158, 94, 155)

                hyperlink = RGBColor(88, 193, 186)
                followed_hyperlink = RGBColor(157, 208, 203)
        #more cases for other themes..
        
# Not working error
#         self.presentation.slide_master.slides.theme.colorscheme[0].rgb = dark1
#         self.presentation.slide_master.theme.colorscheme[1].rgb = light1
#         self.presentation.slide_master.theme.colorscheme[2].rgb = dark2
#         self.presentation.slide_master.theme.colorscheme[3].rgb = light2

#         self.presentation.slide_master.theme.colorscheme[4].rgb = accent1
#         self.presentation.slide_master.theme.colorscheme[5].rgb = accent2
#         self.presentation.slide_master.theme.colorscheme[6].rgb = accent3
#         self.presentation.slide_master.theme.colorscheme[7].rgb = accent4
#         self.presentation.slide_master.theme.colorscheme[8].rgb = accent5
#         self.presentation.slide_master.theme.colorscheme[9].rgb = accent6

#         self.presentation.slide_master.theme.colorscheme[10].rgb = hyperlink
#         self.presentation.slide_master.theme.colorscheme[11].rgb = followed_hyperlink

        # self.set_font_style(title_slide_title.text_frame, font_size=Pt(42), bold=True, italic=False)
#         self.set_font_style(title_slide_content.text_frame, font_size=Pt(20), bold=False, italic=False)
    def addTitleSlide(self,title,subtitle):
        # Title slide 
        title_slide = self.presentation.slides.add_slide(self.title_slide_layout) 
        title_slide_title = title_slide.shapes.title
        title_slide_subtitle = title_slide.placeholders[1]
        title_slide_title.text = str(title)
        title_slide_subtitle.text = str(subtitle)
        
    def addContentSlide(self, title, content):
        # Content slide
        content_slide = self.presentation.slides.add_slide(self.content_slide_layout)
        content_slide_title = content_slide.shapes.title
        content_slide_content = content_slide.placeholders[1]
        content_slide_title.text = str(title)
        content_slide_content.text = str(content)
        
    def savePresentation(self, filename):
        self.setTheme(0)
        self.presentation.save(filename)
        print("Saved Successfully")
    
def CreatePPT(topic,data):
    ppt = PPT(topic)
    for title, content in data.items():
        for index,item in enumerate(content):
            if index == 0:
                ppt.addContentSlide(title, content[index])
            else:
                ppt.addContentSlide(" ", content[index])                
    print("Compilation successful")
    try:
        ppt.savePresentation(f'{topic}.pptx')
    except Exception as e:
        print(e)